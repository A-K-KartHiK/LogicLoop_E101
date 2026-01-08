import streamlit as st
import json
import re
import random
import time
import hashlib
import io
import os
import tempfile
import base64
from pathlib import Path
import subprocess
import sys
import shutil
import atexit

# =========================================================
# 0. SAFE IMPORTS & CONFIG
# =========================================================
try:
    from groq import Groq
    import PyPDF2
    import docx
    from pptx import Presentation
    import graphviz
    
    # Video/Audio imports
    try:
        from PIL import Image, ImageDraw, ImageFont
        import numpy as np
        PIL_AVAILABLE = True
    except ImportError:
        PIL_AVAILABLE = False
        st.warning("PIL/Pillow not available for image generation")
        
    try:
        from gtts import gTTS
        GTTS_AVAILABLE = True
    except ImportError:
        GTTS_AVAILABLE = False
        st.warning("gTTS not available for audio generation")

    # --- MOVIEPY FOR SYNCED VIDEO ---
    try:
        from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
        MOVIEPY_AVAILABLE = True
    except ImportError:
        MOVIEPY_AVAILABLE = False
        st.warning("⚠️ MoviePy is not installed. Video generation will fail. Run: pip install moviepy")
        
except ImportError as e:
    st.error(f"🚨 Required libraries missing! Error: {e}")
    st.info("""
    Run these commands in your terminal to fix everything:
    pip install streamlit groq PyPDF2 python-docx python-pptx graphviz
    pip install moviepy gtts pillow numpy
    """)
    st.stop()

st.set_page_config(page_title="SyllabusQuest: Master Edition", page_icon="🧬", layout="wide")

# CSS STYLING
st.markdown("""
<style>
    .stApp { background-color: #f8f9fa; color: #212529; }
    div.stButton > button {
        width: 100%; border-radius: 8px; font-weight: bold; border: none; padding: 12px;
        transition: 0.3s; background-color: #007bff; color: white;
    }
    div.stButton > button:hover { background-color: #0056b3; transform: scale(1.02); }
    .xp-card {
        background: linear-gradient(135deg, #FFD700, #FFA500); padding: 15px; border-radius: 12px;
        color: black; font-weight: 800; text-align: center; box-shadow: 0 4px 6px rgba(0,0,0,0.1);
    }
    .trust-shield {
        background-color: #d4edda; border-left: 6px solid #28a745; padding: 15px;
        margin-top: 10px; border-radius: 5px; color: #155724; font-size: 0.9rem;
    }
    .game-card {
        background-color: #ffffff; border: 2px solid #e9ecef; padding: 20px; border-radius: 15px;
        box-shadow: 0 4px 10px rgba(0,0,0,0.05); margin-bottom: 20px;
    }
    .video-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 20px; border-radius: 15px; color: white; margin: 10px 0;
    }
    .video-container {
        background: #000; padding: 15px; border-radius: 10px; margin: 10px 0;
    }
</style>
""", unsafe_allow_html=True)

# =========================================================
# 1. API SETUP (USER INPUT ENABLED)
# =========================================================
with st.sidebar:
    st.title("🔐 API Configuration")
    user_api_key = st.text_input("Enter GROQ API Key", type="password", help="Get it from console.groq.com")
    
    if not user_api_key:
        st.warning("⚠️ Enter API Key to start the System")
        st.stop() # Stop execution until key is provided

# Initialize Client with User Key
client = Groq(api_key=user_api_key)

# =========================================================
# 2. VIDEO GENERATION MODULE
# =========================================================
class SimpleVideoGenerator:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
        
    def create_audio_from_text(self, text, topic="topic"):
        """Create audio narration from text"""
        if not GTTS_AVAILABLE:
            return None
        try:
            # Limit text length for TTS to prevent errors
            clean_text = text[:800] if len(text) > 800 else text
            if not clean_text.strip():
                clean_text = f"Educational content about {topic}"
                
            tts = gTTS(text=clean_text, lang='en', slow=False)
            audio_filename = f"audio_{hashlib.md5(topic.encode()).hexdigest()[:8]}.mp3"
            audio_path = os.path.join(self.temp_dir, audio_filename)
            tts.save(audio_path)
            return audio_path
        except Exception as e:
            st.warning(f"Audio creation failed: {e}")
            return None
    
    def create_educational_image(self, topic, content, step=1, total_steps=5):
        """Create educational image for the topic"""
        if not PIL_AVAILABLE:
            return None
        try:
            # Create high-res image
            img = Image.new('RGB', (1280, 720), color=(40, 40, 80))
            draw = ImageDraw.Draw(img)
            
            # Fonts
            try:
                font_large = ImageFont.truetype("arial.ttf", 60)
                font_medium = ImageFont.truetype("arial.ttf", 40)
                font_small = ImageFont.truetype("arial.ttf", 30)
            except:
                font_large = ImageFont.load_default()
                font_medium = ImageFont.load_default()
                font_small = ImageFont.load_default()
            
            # Header
            draw.rectangle([0, 0, 1280, 100], fill=(30, 60, 120))
            draw.text((30, 30), f"🎓 Lesson: {topic}", fill=(255, 255, 255), font=font_medium)
            
            # Progress bar
            progress_width = (step / total_steps) * 1100
            draw.rectangle([50, 120, 50 + progress_width, 140], fill=(0, 200, 100))
            
            # Content Box
            draw.rectangle([50, 180, 1230, 650], outline=(100, 100, 200), width=4, fill=(60, 60, 90))
            
            # Text Wrapping
            words = content.split()
            lines = []
            current_line = []
            
            for word in words:
                current_line.append(word)
                if len(' '.join(current_line)) > 55:
                    lines.append(' '.join(current_line[:-1]))
                    current_line = [word]
            
            if current_line:
                lines.append(' '.join(current_line))
            
            # Draw Text lines
            for i, line in enumerate(lines[:12]):  # Max lines
                y_pos = 220 + i * 45
                draw.text((80, y_pos), line, fill=(240, 240, 255), font=font_small)
            
            # Footer
            draw.text((50, 670), "Generated by SyllabusQuest AI", fill=(100, 100, 150), font=font_small)
            
            # Save
            img_filename = f"frame_{hashlib.md5((topic+str(step)).encode()).hexdigest()[:8]}.png"
            img_path = os.path.join(self.temp_dir, img_filename)
            img.save(img_path)
            
            return img_path
            
        except Exception as e:
            st.error(f"Image creation error: {e}")
            return None

    def create_video_frames(self, topic, content):
        """Create multiple frames for the video"""
        if not PIL_AVAILABLE: return []
        try:
            frames = []
            # Intro Frame
            frames.append(self.create_educational_image(topic, f"Introduction to {topic}", 1, 5))
            
            # Content Frames (Split content)
            words = content.split()
            chunk_size = 40
            chunks = [' '.join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]
            
            for i, chunk in enumerate(chunks[:3]): # Max 3 content slides
                frames.append(self.create_educational_image(topic, chunk, i+2, 5))
                
            # Summary Frame
            frames.append(self.create_educational_image(topic, "Key Takeaways & Summary", 5, 5))
            
            return [f for f in frames if f is not None]
        except Exception as e:
            st.error(f"Frame creation error: {e}")
            return []

    def render_final_video(self, frames, audio_path, output_filename="output.mp4"):
        """
        THE KEY FIX: Uses MoviePy to stitch images and audio into a real MP4.
        This ensures perfect sync and playback support.
        """
        if not MOVIEPY_AVAILABLE:
            st.error("MoviePy not found. Cannot render video file.")
            return None
            
        try:
            output_path = os.path.join(self.temp_dir, output_filename)
            
            # 1. Load Audio to get duration
            audio_clip = AudioFileClip(audio_path)
            total_duration = audio_clip.duration
            
            # 2. Calculate duration per frame
            if not frames: return None
            duration_per_frame = total_duration / len(frames)
            
            # 3. Create Image Clips
            clips = []
            for frame_path in frames:
                clip = ImageClip(frame_path).set_duration(duration_per_frame)
                clips.append(clip)
            
            # 4. Concatenate
            video = concatenate_videoclips(clips, method="compose")
            
            # 5. Set Audio
            video = video.set_audio(audio_clip)
            
            # 6. Write File (using codec compatible with browsers)
            # preset='ultrafast' creates video quickly
            video.write_videofile(
                output_path, 
                fps=24, 
                codec='libx264', 
                audio_codec='aac', 
                preset='ultrafast',
                logger=None # Suppress terminal output
            )
            
            return output_path
            
        except Exception as e:
            st.error(f"Rendering failed: {e}")
            return None

    def cleanup(self):
        """Clean up temporary files"""
        try:
            shutil.rmtree(self.temp_dir, ignore_errors=True)
        except:
            pass

# Initialize generator
video_gen = SimpleVideoGenerator()
atexit.register(video_gen.cleanup)

# =========================================================
# 3. INTELLIGENT ENGINE (Optimized & Diagram Aware)
# =========================================================
def get_groq_response(prompt, context_text, expect_json=False, temperature=0.3):
    """Core AI Engine. Handles Chunking, JSON enforcement, and Context."""
    try:
        # Reduced context size slightly to prevent timeouts/errors
        safe_context = context_text[:15000] 
        
        full_prompt = f"""
You are an academic assistant. Use ONLY the following context as the knowledge source.
If something is not in the context, say you don't know and NEVER invent facts.

CONTEXT (Uploaded syllabus / notes):
{safe_context}

TASK:
{prompt}
"""
        # Improved System Prompt for Diagrams
        system_msg = """
You are a Learning-Aware AI. Adapt to Beginner/Intermediate/Advanced levels.
Assess if the user would understand the response better with a diagram.
If yes, insert a diagram tag 

[Image of X]
 where X is a specific, contextually relevant query.
Place the tag immediately before or after relevant text.
Do NOT use tags for generic illustrations.
"""

        if expect_json:
            full_prompt += "\n\nCRITICAL: Return ONLY valid JSON. No Markdown. No Intro."

        completion = client.chat.completions.create(
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": full_prompt}
            ],
            model="llama-3.1-8b-instant", # Using fast model to avoid rate limits
            temperature=temperature,
            response_format={"type": "json_object"} if expect_json else None
        )

        response_text = completion.choices[0].message.content

        if expect_json:
            try:
                # Robust extraction: find first { and last }
                clean_text = re.sub(r"```json|```", "", response_text).strip()
                start = clean_text.find('{')
                end = clean_text.rfind('}') + 1
                if start != -1 and end != -1:
                    return json.loads(clean_text[start:end])
                else:
                    return json.loads(clean_text) # Try raw
            except Exception as json_err:
                print(f"JSON Parsing Error: {json_err}")
                return None

        return response_text
    except Exception as e:
        # VISIBLE ERROR MESSAGE FOR DEBUGGING
        if "429" in str(e):
             st.error("🚨 Rate Limit Reached. Please wait a moment before trying again.")
        else:
             st.error(f"🚨 AI Error: {e}")
        return None

# =========================================================
# 4. HELPERS
# =========================================================
def get_topic_image(topic):
    """Dynamically selects an image based on topic keywords."""
    t = str(topic).lower()
    base = "https://images.unsplash.com/"
    if any(x in t for x in ["cell", "bio", "dna"]): return base + "photo-1530026405186-ed1f139313f8?w=800"
    if any(x in t for x in ["space", "star"]): return base + "photo-1451187580459-43490279c0fa?w=800"
    if any(x in t for x in ["atom", "force"]): return base + "photo-1635070041078-e363dbe005cb?w=800"
    if any(x in t for x in ["code", "computer"]): return base + "photo-1555066931-4365d14bab8c?w=800"
    return base + "photo-1456513080510-7bf3a84b82f8?w=800"

@st.cache_data(show_spinner=False)
def extract_file_content(uploaded_file):
    """Extracts text from PDF, DOCX, PPTX, TXT."""
    text_content = ""
    try:
        if uploaded_file.name.endswith(".pdf"):
            pdf = PyPDF2.PdfReader(uploaded_file)
            for page in pdf.pages:
                text_content += page.extract_text() or ""
        elif uploaded_file.name.endswith(".docx"):
            doc = docx.Document(uploaded_file)
            text_content = "\n".join([para.text for para in doc.paragraphs])
        elif uploaded_file.name.endswith(".pptx"):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif uploaded_file.name.endswith(".txt"):
            return str(uploaded_file.read(), "utf-8")
    except Exception:
        return None
    return text_content

# =========================================================
# 5. SESSION STATE INIT
# =========================================================
if 'file_text' not in st.session_state: st.session_state.file_text = ""
if 'syllabus' not in st.session_state: st.session_state.syllabus = []
if 'current_topic_index' not in st.session_state: st.session_state.current_topic_index = 0
if 'xp' not in st.session_state: st.session_state.xp = 0
if 'total_qs' not in st.session_state: st.session_state.total_qs = 0
if 'correct_qs' not in st.session_state: st.session_state.correct_qs = 0
if 'lesson_content' not in st.session_state: st.session_state.lesson_content = None
if 'quiz_card' not in st.session_state: st.session_state.quiz_card = None
if 'exam_paper' not in st.session_state: st.session_state.exam_paper = None
if 'exam_answers' not in st.session_state: st.session_state.exam_answers = {}
if 'chat_history' not in st.session_state: st.session_state.chat_history = []
if 'card_revealed' not in st.session_state: st.session_state.card_revealed = False

# New Video States
if 'generated_videos' not in st.session_state: st.session_state.generated_videos = {} # Store topic:path

# =========================================================
# 6. SIDEBAR & UPLOAD
# =========================================================
with st.sidebar:
    st.title("📂 Knowledge Base")
    uploaded_file = st.file_uploader("Upload File", type=['pdf', 'docx', 'pptx', 'txt'])

    if uploaded_file and not st.session_state.file_text:
        with st.spinner("🧠 Analyzing & Creating Syllabus..."):
            text = extract_file_content(uploaded_file)
            if text:
                st.session_state.file_text = text
                syl_prompt = (
                    "From the context, list the top 5-8 main academic concepts/chapters ONLY.\n"
                    "JSON format: {\"topics\": [\"Topic 1\", \"Topic 2\", ...]}"
                )
                syl_data = get_groq_response(syl_prompt, text, expect_json=True)

                if syl_data and 'topics' in syl_data:
                    st.session_state.syllabus = syl_data['topics']
                    st.success("✅ File Indexed!")
                else:
                    st.session_state.syllabus = ["General Content"]
            else:
                st.error("❌ Could not read file.")

    # NAVIGATION
    if st.session_state.syllabus:
        st.subheader("📖 Syllabus Navigator")
        selected_topic = st.selectbox(
            "Jump to Concept:",
            st.session_state.syllabus,
            index=st.session_state.current_topic_index
        )
        if selected_topic != st.session_state.syllabus[st.session_state.current_topic_index]:
            st.session_state.current_topic_index = st.session_state.syllabus.index(selected_topic)
            st.rerun()

    st.divider()
    st.markdown("### 📊 Performance")
    st.markdown(f"<div class='xp-card'>{st.session_state.xp} XP</div>", unsafe_allow_html=True)

# =========================================================
# 7. MAIN APP
# =========================================================
st.title("🧬 SyllabusQuest: Master Edition")

if st.session_state.file_text:

    tabs = st.tabs(["📚 Adaptive Lesson", "🎮 Endless Game", "⚔️ Interactive Exam", 
                    "⚡ 1-Hour Revision", "📈 Analytics", "💬 Neural Chat", 
                    "🎥 Video Studio", "⚖️ Safety Audit"])

    # ---------------------------------------------------------
    # TAB 1: ADAPTIVE LESSON
    # ---------------------------------------------------------
    with tabs[0]:
        c1, c2 = st.columns([1, 3])
        with c1:
            st.subheader("Settings")
            current_topic = st.session_state.syllabus[st.session_state.current_topic_index]
            st.info(f"Current: **{current_topic}**")
            lvl = st.radio("Level", ["Beginner", "Intermediate", "Advanced"], key="l1")
            style = st.radio("Style", ["Visual", "Real-World", "Academic"], key="s1")

            if st.button("🚀 Teach This"):
                with st.spinner("Generating..."):
                    prompt = f"""
                    Teach '{current_topic}' based ONLY on the context. 
                    Level: {lvl}. Style: {style}.
                    
                    REQUIREMENT: Generate EXACTLY 5 distinct Multiple Choice Questions (quiz).
                    
                    Output JSON: {{
                        "title": "Lesson Title", 
                        "content": "Detailed explanation...", 
                        "real_world": "Real world example...", 
                        "citation": "Source...", 
                        "quiz": [
                            {{"q":"Question 1","opts":["A) ...","B) ...","C) ...", "D) ..."],"ans":"A","reason":"Explanation for Q1"}},
                            {{"q":"Question 2","opts":["A) ...","B) ...","C) ...", "D) ..."],"ans":"B","reason":"Explanation for Q2"}},
                            {{"q":"Question 3","opts":["A) ...","B) ...","C) ...", "D) ..."],"ans":"C","reason":"Explanation for Q3"}},
                            {{"q":"Question 4","opts":["A) ...","B) ...","C) ...", "D) ..."],"ans":"D","reason":"Explanation for Q4"}},
                            {{"q":"Question 5","opts":["A) ...","B) ...","C) ...", "D) ..."],"ans":"A","reason":"Explanation for Q5"}}
                        ]
                    }}
                    """
                    data = get_groq_response(prompt, st.session_state.file_text, expect_json=True)
                    if data: 
                        st.session_state.lesson_content = data
                    else:
                        st.error("⚠️ AI returned no content. Please check API Key or File Content.")

        with c2:
            if st.session_state.lesson_content:
                d = st.session_state.lesson_content
                st.image(get_topic_image(d.get('title', current_topic)), use_container_width=True)
                st.markdown(f"## {d.get('title')}")
                st.write(d.get('content'))
                if d.get('real_world'):
                    st.info(f"🌍 **Real World:** {d['real_world']}")
                
                st.markdown("---")
                st.subheader(f"🧠 {lvl} Quiz (5 Questions)")
                
                # Render 5 Questions
                quiz_data = d.get('quiz', [])
                for i, q in enumerate(quiz_data):
                    with st.expander(f"Question {i+1}: {q.get('q')}", expanded=True):
                        cols = st.columns(4)
                        for idx, opt in enumerate(q.get('opts', [])):
                            # Unique key for every button
                            if cols[idx].button(opt, key=f"q_{i}_{idx}", use_container_width=True):
                                label = opt.split(")")[0].strip()
                                correct_ans = q.get('ans').strip()
                                
                                if label == correct_ans:
                                    st.success(f"✅ Correct! {q.get('reason')}")
                                    st.session_state.xp += 10
                                    st.session_state.correct_qs += 1
                                else:
                                    st.error(f"❌ Wrong. {q.get('reason')}")
                                st.session_state.total_qs += 1

    # ---------------------------------------------------------
    # TAB 2: ENDLESS GAME
    # ---------------------------------------------------------
    with tabs[1]:
        st.subheader("🎮 Knowledge Arena")
        c_game, c_ctrl = st.columns([3, 1])

        with c_ctrl:
            # Main control for new card
            if not st.session_state.quiz_card:
                if st.button("🎲 Deal First Card"):
                    topic_card = random.choice(st.session_state.syllabus)
                    prompt = f"""Create 1 MCQ for '{topic_card}'. JSON: {{"q":"...","opts":["A)...","B)...","C)...","D)..."],"ans":"A","exp":"..."}}"""
                    st.session_state.quiz_card = get_groq_response(prompt, st.session_state.file_text, expect_json=True)
                    st.session_state.card_revealed = False
                    st.rerun()
            else:
                 # Skip logic
                if st.button("⏩ Skip / Show Answer"):
                    st.session_state.card_revealed = True
                    st.rerun()

                # Next Question Logic
                if st.button("⏭️ Next Question"):
                    st.session_state.quiz_card = None
                    st.session_state.card_revealed = False
                    
                    # Immediately generate new one
                    topic_card = random.choice(st.session_state.syllabus)
                    prompt = f"""Create 1 MCQ for '{topic_card}'. JSON: {{"q":"...","opts":["A)...","B)...","C)...","D)..."],"ans":"A","exp":"..."}}"""
                    st.session_state.quiz_card = get_groq_response(prompt, st.session_state.file_text, expect_json=True)
                    st.rerun()

        with c_game:
            if st.session_state.quiz_card:
                q = st.session_state.quiz_card
                st.markdown(f"<div class='game-card'><h3>❓ {q.get('q')}</h3></div>", unsafe_allow_html=True)
                
                cols = st.columns(2)
                if not st.session_state.card_revealed:
                    for i, opt in enumerate(q.get('opts')):
                        if cols[i%2].button(opt, key=f"gm_{i}", use_container_width=True):
                            st.session_state.card_revealed = True
                            label = opt.split(")")[0].strip()
                            if label == q.get('ans'):
                                st.balloons()
                                st.success("🏆 CORRECT! +50 XP")
                                st.session_state.xp += 50
                            else:
                                st.error(f"💀 WRONG! Answer: {q.get('ans')}")
                            st.rerun()
                else:
                    st.info(f"**Correct Answer:** {q.get('ans')}")
                    st.warning(f"**Explanation:** {q.get('exp')}")

    # ---------------------------------------------------------
    # TAB 3: INTERACTIVE EXAM (UPDATED AS REQUESTED)
    # ---------------------------------------------------------
    with tabs[2]:
        st.subheader("⚔️ Examination Hall")
        # CHANGED: ONLY MCQ AND FILL IN THE BLANKS
        q_type = st.selectbox("Type:", ["MCQ", "Fill in the Blanks"])
        diff = st.radio("Difficulty:", ["Easy", "Medium", "Hard"], horizontal=True)

        if st.button("📄 Generate Exam"):
            with st.spinner("Setting Paper..."):
                # STRICT PROMPT FOR ONLY MCQ (4 OPTIONS) AND FILL IN THE BLANKS
                prompt = f"""
                Create a strict JSON object with exactly 5 {q_type} questions based on the text. Difficulty: {diff}.
                
                JSON Structure must be EXACTLY:
                {{
                  "questions": [
                    {{
                      "id": 1,
                      "type": "{q_type}",
                      "text": "Question text here (For Fill in Blanks, use '_______' for the blank)",
                      "options": ["A) ...", "B) ...", "C) ...", "D) ..."] (ONLY IF MCQ),
                      "correct": "Answer here"
                    }}
                  ]
                }}

                IMPORTANT:
                1. If type is MCQ, 'options' must be a list of 4 distinct strings.
                2. If type is Fill in the Blanks, 'options' must be an empty list [].
                3. Ensure valid JSON syntax (close all brackets/braces).
                """
                exam_data = get_groq_response(prompt, st.session_state.file_text, expect_json=True)
                if exam_data:
                    st.session_state.exam_paper = exam_data.get('questions', [])
                    st.session_state.exam_answers = {}
                else:
                    st.error("⚠️ Failed to generate exam. Please try again.")

        if st.session_state.exam_paper:
            with st.form("exam_form"):
                for q in st.session_state.exam_paper:
                    st.write(f"**{q['id']}. {q['text']}**")
                    
                    if q_type == "MCQ":
                        # Ensure 4 options exist
                        options = q.get('options', ["Error: No options generated"])
                        st.session_state.exam_answers[q['id']] = st.radio(
                            "Select Answer:", options, key=f"ex_{q['id']}"
                        )
                    else:
                        # Fill in the blanks logic
                        st.session_state.exam_answers[q['id']] = st.text_input("Type your answer here:", key=f"ex_{q['id']}")
                    st.markdown("---")
                
                if st.form_submit_button("Submit Exam"):
                    score = 0
                    for q in st.session_state.exam_paper:
                        ans = st.session_state.exam_answers.get(q['id'])
                        
                        if q_type == "MCQ":
                            if ans and ans.startswith(q.get('correct', 'X')):
                                score += 1
                                st.success(f"Q{q['id']}: Correct")
                            else:
                                st.error(f"Q{q['id']}: Wrong. Correct: {q.get('correct')}")
                        else:
                            # Fill in the blanks strict grading
                            correct_val = q.get('correct', '').strip().lower()
                            user_val = str(ans).strip().lower() if ans else ""
                            if user_val == correct_val:
                                score += 1
                                st.success(f"Q{q['id']}: Correct")
                            else:
                                st.error(f"Q{q['id']}: Wrong. Correct: {q.get('correct')}")
                    
                    st.metric("Score", f"{score}/5")
                    if score == 5: st.session_state.xp += 100

    # ---------------------------------------------------------
    # TAB 4: REVISION
    # ---------------------------------------------------------
    with tabs[3]:
        st.subheader("⚡ 1-Hour Revision")
        if st.button("🔥 Generate Notes"):
            with st.spinner("Analysing Context & Summarizing..."):
                prompt = """
                Based ONLY on the provided context, create a revision sheet with:
                1. 5 Key Definitions
                2. 3 Common Misconceptions
                3. A Formula/Date Cheat Sheet
                4. A Golden Summary
                Format in clean Markdown.
                """
                # expect_json=False because we want Markdown
                rev = get_groq_response(prompt, st.session_state.file_text, expect_json=False)
                if rev:
                    st.markdown(rev)
                else:
                    st.error("⚠️ AI Error. Please check your API key.")

    # ---------------------------------------------------------
    # TAB 5: ANALYTICS
    # ---------------------------------------------------------
    with tabs[4]:
        st.subheader("📈 Analytics")
        c1, c2 = st.columns(2)
        c1.metric("Total XP", st.session_state.xp)
        c2.metric("Questions Done", st.session_state.total_qs)
        st.bar_chart({"Correct": st.session_state.correct_qs, "Wrong": st.session_state.total_qs - st.session_state.correct_qs})

    # ---------------------------------------------------------
    # TAB 6: NEURAL CHAT
    # ---------------------------------------------------------
    with tabs[5]:
        st.subheader("💬 AI Tutor")
        for m in st.session_state.chat_history:
            with st.chat_message(m["role"]): st.write(m["content"])
        
        if p := st.chat_input("Ask a specific doubt..."):
            st.session_state.chat_history.append({"role": "user", "content": p})
            st.chat_message("user").write(p)
            
            # Simple direct prompt
            r = get_groq_response(f"Answer this question strictly using the provided context: {p}", st.session_state.file_text, expect_json=False)
            
            if not r: r = "⚠️ Error: I could not reach the AI service. Please check your API key."
            
            st.session_state.chat_history.append({"role": "assistant", "content": r})
            st.chat_message("assistant").write(r)

    # ---------------------------------------------------------
    # TAB 7: VIDEO STUDIO (UNTOUCHED)
    # ---------------------------------------------------------
    with tabs[6]:
        st.subheader("🎥 Video Studio")
        st.markdown("""
        **Note:** This module now generates a real `.mp4` file to ensure perfect audio-video synchronization. 
        Please allow 30-60 seconds for rendering.
        """)

        c_gen, c_play = st.columns([1, 2])

        with c_gen:
            st.markdown("### ⚙️ Generator")
            v_topic = st.selectbox("Select Topic", st.session_state.syllabus, key="vid_top")
            
            if st.button("🎬 Render Video", type="primary"):
                with st.spinner("1/4 Writing Script..."):
                    content = get_groq_response(f"Explain '{v_topic}' for a video. Plain text only.", st.session_state.file_text)
                
                with st.spinner("2/4 Generating Audio..."):
                    audio_path = video_gen.create_audio_from_text(content, v_topic)
                
                with st.spinner("3/4 Generating Frames..."):
                    frames = video_gen.create_video_frames(v_topic, content)
                
                with st.spinner("4/4 Rendering Final MP4 (Syncing)..."):
                    if frames and audio_path:
                        output_file = f"video_{hashlib.md5(v_topic.encode()).hexdigest()[:5]}.mp4"
                        video_path = video_gen.render_final_video(frames, audio_path, output_file)
                        
                        if video_path:
                            st.session_state.generated_videos[v_topic] = video_path
                            st.success("✅ Video Rendered!")
                        else:
                            st.error("❌ Rendering failed. Check MoviePy installation.")
                    else:
                        st.error("❌ Failed to generate assets.")

        with c_play:
            st.markdown("### 📺 Player")
            if st.session_state.generated_videos:
                play_topic = st.selectbox("Select Video to Play", list(st.session_state.generated_videos.keys()))
                video_file = st.session_state.generated_videos[play_topic]
                
                if os.path.exists(video_file):
                    # Using native Streamlit video player for perfect playback
                    st.video(video_file)
                    
                    with open(video_file, "rb") as f:
                        st.download_button("📥 Download MP4", f, file_name=f"{play_topic}.mp4")
                else:
                    st.error("Video file not found. Please regenerate.")
            else:
                st.info("No videos generated yet. Use the panel on the left.")

    # ---------------------------------------------------------
    # TAB 8: SAFETY AUDIT
    # ---------------------------------------------------------
    with tabs[7]:
        st.subheader("⚖️ Safety Audit")
        if st.button("🔍 Run Audit"):
            with st.spinner("Auditing..."):
                audit = get_groq_response("Audit the provided content for any hallucinations or educational bias.", st.session_state.file_text, expect_json=False)
                if audit:
                    st.markdown(audit)
                    st.success("Audit Complete")
                else:
                    st.error("Audit failed to generate.")

else:
    st.info("👆 Upload a file to begin.")

# Cleanup hook
import atexit
atexit.register(video_gen.cleanup)